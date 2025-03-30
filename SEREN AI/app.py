from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
from youtube_transcript_api import YouTubeTranscriptApi
import google.generativeai as genai
import tempfile
import os
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from pptx import Presentation
import docx


app = Flask(__name__, template_folder='templates')
CORS(app)  # Enable CORS for all routes
# Configure API keys and secret key
app.config['YOUTUBE_API_KEY'] = os.getenv('YOUTUBE_API_KEY')
app.config['GEMINI_API_KEY'] = os.getenv('GEMINI_API_KEY')
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY')

app = Flask(__name__, template_folder='templates')
CORS(app)  # Enable CORS for all routes

# Configuration
app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB limit
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'txt', 'doc', 'docx'}

# Initialize Gemini
genai.configure(api_key='')

@app.route('/')
def home():
    return render_template('Index.html')

@app.route('/get_transcript', methods=['GET'])
def get_transcript():
    video_id = request.args.get('video_id')
    if not video_id:
        return jsonify({"error": "video_id parameter is required"}), 400
    
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join([entry['text'] for entry in transcript])
        return jsonify({"transcript": transcript_text})
    except Exception as e:
        return jsonify({"error": f"Failed to get transcript: {str(e)}"}), 500

@app.route('/generate_summary', methods=['POST'])
def generate_summary():
    if not request.is_json:
        return jsonify({"error": "Request must be JSON"}), 400
        
    data = request.get_json()
    transcript = data.get('transcript', '')
    
    if not transcript:
        return jsonify({"error": "No transcript provided"}), 400
    
    try:
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        prompt = f"Generate a concise summary (about 100 words) of the following video transcript:\n\n{transcript}"
        response = model.generate_content(prompt)
        
        if not response.text:
            return jsonify({"error": "Failed to generate summary"}), 500
            
        return jsonify({"summary": response.text})
    except Exception as e:
        return jsonify({"error": f"Summary generation failed: {str(e)}"}), 500

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(filepath):
    try:
        if filepath.lower().endswith('.pdf'):
            with open(filepath, 'rb') as f:
                return "\n".join(page.extract_text() for page in PdfReader(f).pages)
        elif filepath.lower().endswith(('.ppt', '.pptx')):
            return "\n".join(shape.text for slide in Presentation(filepath).slides 
                          for shape in slide.shapes if hasattr(shape, "text"))
        elif filepath.lower().endswith(('.doc', '.docx')):
            return "\n".join(p.text for p in docx.Document(filepath).paragraphs)
        elif filepath.lower().endswith('.txt'):
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        raise ValueError("Unsupported file format")
    except Exception as e:
        raise ValueError(f"Error extracting text: {str(e)}")

def generate_quiz(text):
    try:
        prompt = f"""Generate 5 clear multiple-choice questions based on the following text.
For each question, provide 4 options (A-D) and indicate the correct answer.

Text:
{text}

Format each question exactly like this:
Question: [question text]
A) Option 1
B) Option 2
C) Option 3
D) Option 4
Answer: [correct letter]"""
        
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("No quiz generated")
            
        return response.text
    except Exception as e:
        raise ValueError(f"Quiz generation failed: {str(e)}")

@app.route('/api/generate-quiz-from-file', methods=['POST', 'OPTIONS'])
def handle_file_upload():
    if request.method == 'OPTIONS':
        response = jsonify({'status': 'preflight'})
        response.headers.add('Access-Control-Allow-Origin', '*')
        response.headers.add('Access-Control-Allow-Headers', '*')
        response.headers.add('Access-Control-Allow-Methods', 'POST, OPTIONS')
        return response
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'File type not allowed. Allowed types: ' + ', '.join(ALLOWED_EXTENSIONS)}), 400

    filepath = None
    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        text = extract_text_from_file(filepath)
        quiz = generate_quiz(text)
        
        return jsonify({
            'quiz': quiz,
            'filename': filename,
            'text_length': len(text)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        if filepath and os.path.exists(filepath):
            os.remove(filepath)

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)